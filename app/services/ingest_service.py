import os
import base64
import zipfile
from io import BytesIO

import fitz  # pymupdf
from PIL import Image
from groq import Groq
from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE
from langchain_community.embeddings import HuggingFaceEmbeddings
from langchain_community.vectorstores import FAISS
from langchain_text_splitters import RecursiveCharacterTextSplitter
from langchain_core.documents import Document

from app.core.config import EMBEDDING_MODEL, GROQ_API_KEY, SUPABASE_URL, SUPABASE_KEY

VECTOR_BASE_PATH = "vectorstore"

# ---------------------------------------------------------------------------
# Clients (lazy init)
# ---------------------------------------------------------------------------

_groq_client = None
_supabase_client = None


def get_groq():
    global _groq_client
    if _groq_client is None:
        _groq_client = Groq(api_key=GROQ_API_KEY)
    return _groq_client


def get_supabase():
    global _supabase_client
    if _supabase_client is None:
        from supabase import create_client
        _supabase_client = create_client(SUPABASE_URL, SUPABASE_KEY)
    return _supabase_client


# ---------------------------------------------------------------------------
# Supabase Storage — upload raw file
# ---------------------------------------------------------------------------

def upload_to_storage(file_path: str, subject_id: str, original_name: str) -> str:
    """Upload raw file to Supabase Storage. Returns public URL."""
    bucket = "study-notes"
    storage_path = f"{subject_id}/{original_name}"

    with open(file_path, "rb") as f:
        file_bytes = f.read()

    mime = "application/pdf"
    if original_name.endswith(".pptx"):
        mime = "application/vnd.openxmlformats-officedocument.presentationml.presentation"
    elif original_name.endswith(".docx"):
        mime = "application/vnd.openxmlformats-officedocument.wordprocessingml.document"
    elif original_name.lower().endswith((".png", ".jpg", ".jpeg")):
        mime = "image/png"

    try:
        sb = get_supabase()
        sb.storage.from_(bucket).upload(
            storage_path,
            file_bytes,
            {"content-type": mime, "upsert": "true"}
        )
        url = sb.storage.from_(bucket).get_public_url(storage_path)
        print(f"  [storage] Uploaded {original_name} → {url}")
        return url
    except Exception as e:
        print(f"  [storage] Upload failed for {original_name}: {e}")
        return ""


# ---------------------------------------------------------------------------
# Equation OCR via pix2tex (LaTeX)
# ---------------------------------------------------------------------------

def image_to_latex(image_bytes: bytes) -> str:
    """
    Try to extract a LaTeX formula from an image.
    Returns empty string if pix2tex is unavailable or image is not an equation.
    """
    try:
        from pix2tex.cli import LatexOCR
        img = Image.open(BytesIO(image_bytes)).convert("RGB")
        model = LatexOCR()
        latex = model(img).strip()
        # Only return if it looks like actual math (contains letters + operators)
        if latex and any(c in latex for c in ["\\", "=", "^", "_", "frac"]):
            return f"$$ {latex} $$"
    except ImportError:
        pass  # pix2tex not installed — skip silently
    except Exception:
        pass
    return ""


# ---------------------------------------------------------------------------
# Vision: describe image via Groq
# ---------------------------------------------------------------------------

def describe_image(image_bytes: bytes, surrounding_context: str = "", source_hint: str = "") -> str:
    """Describe an image using Groq vision. Returns empty string for tiny/blank images."""
    try:
        img = Image.open(BytesIO(image_bytes))
        if img.width < 80 or img.height < 80:
            return ""

        # Try equation OCR first — more precise than a vision description for math
        latex = image_to_latex(image_bytes)
        if latex:
            return f"Mathematical equation extracted from figure: {latex}"

        buf = BytesIO()
        img.convert("RGB").save(buf, format="JPEG", quality=85)
        b64 = base64.b64encode(buf.getvalue()).decode()

        ctx = f"\nSurrounding text: {surrounding_context[:300]}" if surrounding_context else ""
        src = f"\nSource: {source_hint}" if source_hint else ""

        response = get_groq().chat.completions.create(
            model="meta-llama/llama-4-scout-17b-16e-instruct",
            messages=[{
                "role": "user",
                "content": [
                    {"type": "image_url", "image_url": {"url": f"data:image/jpeg;base64,{b64}"}},
                    {"type": "text", "text": (
                        f"You are processing a study note or academic slide.{ctx}{src}\n\n"
                        "Describe this image for a student who cannot see it:\n"
                        "- What type of figure is it? (diagram, graph, chart, flowchart, table, photo, etc.)\n"
                        "- What concept does it illustrate?\n"
                        "- Label ALL axes, components, boxes, arrows, and annotations.\n"
                        "- Extract ALL numerical values, formulas, or equations.\n"
                        "- Describe key trends, relationships, or takeaways.\n"
                        "Be thorough — this description is used for retrieval."
                    )}
                ]
            }],
            max_tokens=600
        )
        return response.choices[0].message.content.strip()

    except Exception as e:
        print(f"  [vision] Could not describe image: {e}")
        return ""


# ---------------------------------------------------------------------------
# File loaders
# ---------------------------------------------------------------------------

def load_pdf(file_path: str, original_name: str) -> list[Document]:
    docs = []
    doc = fitz.open(file_path)

    for page_num, page in enumerate(doc, start=1):
        text = page.get_text("text").strip()
        if text:
            docs.append(Document(
                page_content=text,
                metadata={"source": original_name, "page": page_num, "type": "text"}
            ))

        for img_index, img_info in enumerate(page.get_images(full=True)):
            xref = img_info[0]
            try:
                base_image = doc.extract_image(xref)
                description = describe_image(base_image["image"], surrounding_context=text, source_hint=original_name)
                if description:
                    print(f"  [vision] Page {page_num}, image {img_index+1} — described")
                    docs.append(Document(
                        page_content=f"[Figure on page {page_num}]: {description}",
                        metadata={"source": original_name, "page": page_num, "type": "image_description"}
                    ))
            except Exception as e:
                print(f"  [vision] Skipped image on page {page_num}: {e}")

    doc.close()
    return docs


def load_pptx(file_path: str, original_name: str) -> list[Document]:
    docs = []
    prs = Presentation(file_path)

    for slide_num, slide in enumerate(prs.slides, start=1):
        text_parts = []
        image_blobs = []

        for shape in slide.shapes:
            if hasattr(shape, "text") and shape.text.strip():
                text_parts.append(shape.text.strip())
            if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
                try:
                    image_blobs.append(shape.image.blob)
                except Exception:
                    pass

        slide_text = "\n".join(text_parts).strip()
        if slide_text:
            docs.append(Document(
                page_content=slide_text,
                metadata={"source": original_name, "slide": slide_num, "type": "text"}
            ))

        for idx, blob in enumerate(image_blobs):
            description = describe_image(blob, surrounding_context=slide_text, source_hint=original_name)
            if description:
                print(f"  [vision] Slide {slide_num}, image {idx+1} — described")
                docs.append(Document(
                    page_content=f"[Figure on slide {slide_num}]: {description}",
                    metadata={"source": original_name, "slide": slide_num, "type": "image_description"}
                ))

    # Safety net: scan ppt/media/ in the zip for any missed images
    try:
        with zipfile.ZipFile(file_path) as z:
            for media_path in z.namelist():
                if media_path.startswith("ppt/media/") and media_path.lower().endswith(
                    (".png", ".jpg", ".jpeg", ".gif", ".bmp")
                ):
                    description = describe_image(z.read(media_path), source_hint=original_name)
                    if description:
                        docs.append(Document(
                            page_content=f"[Embedded media — {os.path.basename(media_path)}]: {description}",
                            metadata={"source": original_name, "type": "image_description"}
                        ))
    except Exception as e:
        print(f"  [pptx-zip] Could not scan media: {e}")

    return docs


def load_docx(file_path: str, original_name: str) -> list[Document]:
    from docx import Document as DocxDocument
    docs = []

    try:
        docx = DocxDocument(file_path)
        full_text = "\n".join(p.text for p in docx.paragraphs if p.text.strip())
        if full_text:
            docs.append(Document(
                page_content=full_text,
                metadata={"source": original_name, "type": "text"}
            ))

        with zipfile.ZipFile(file_path) as z:
            for media_path in z.namelist():
                if media_path.startswith("word/media/") and media_path.lower().endswith(
                    (".png", ".jpg", ".jpeg", ".gif", ".bmp")
                ):
                    description = describe_image(
                        z.read(media_path),
                        surrounding_context=full_text[:300],
                        source_hint=original_name
                    )
                    if description:
                        docs.append(Document(
                            page_content=f"[Figure in {original_name}]: {description}",
                            metadata={"source": original_name, "type": "image_description"}
                        ))
    except Exception as e:
        print(f"  [docx] Error loading {original_name}: {e}")

    return docs


# ---------------------------------------------------------------------------
# Chunking
# ---------------------------------------------------------------------------

def split_documents(documents: list[Document]) -> list[Document]:
    splitter = RecursiveCharacterTextSplitter(chunk_size=500, chunk_overlap=100)
    chunks = []
    for doc in documents:
        if doc.metadata.get("type") == "image_description":
            chunks.append(doc)  # never split image descriptions
        else:
            for chunk in splitter.split_text(doc.page_content):
                chunks.append(Document(page_content=chunk, metadata=doc.metadata))
    return chunks


# ---------------------------------------------------------------------------
# FAISS vectorstore (per subject)
# ---------------------------------------------------------------------------

def get_vector_path(subject_id: str) -> str:
    return os.path.join(VECTOR_BASE_PATH, subject_id)


def upsert_vectorstore(chunks: list[Document], subject_id: str):
    embeddings = HuggingFaceEmbeddings(model_name=EMBEDDING_MODEL)
    vector_path = get_vector_path(subject_id)

    if os.path.exists(vector_path):
        db = FAISS.load_local(vector_path, embeddings, allow_dangerous_deserialization=True)
        db.add_documents(chunks)
    else:
        db = FAISS.from_documents(chunks, embeddings)

    os.makedirs(vector_path, exist_ok=True)
    db.save_local(vector_path)
    print(f"  [faiss] Saved {len(chunks)} chunks for subject '{subject_id}'")


# ---------------------------------------------------------------------------
# Public entry point
# ---------------------------------------------------------------------------

def ingest_file(file_path: str, subject_id: str, original_name: str = None):
    if original_name is None:
        original_name = os.path.basename(file_path)

    ext = original_name.lower()

    # 1. Upload raw file to Supabase Storage
    storage_url = upload_to_storage(file_path, subject_id, original_name)

    # 2. Extract text + image descriptions
    if ext.endswith(".pdf"):
        docs = load_pdf(file_path, original_name)
    elif ext.endswith(".pptx"):
        docs = load_pptx(file_path, original_name)
    elif ext.endswith(".docx"):
        docs = load_docx(file_path, original_name)
    elif ext.endswith((".png", ".jpg", ".jpeg")):
        with open(file_path, "rb") as f:
            image_bytes = f.read()
        description = describe_image(image_bytes, source_hint=original_name)
        docs = [Document(
            page_content=f"[Image file — {original_name}]: {description}",
            metadata={"source": original_name, "type": "image_description"}
        )] if description else []
    else:
        raise ValueError(f"Unsupported file type: {original_name}")

    if not docs:
        print(f"  [ingest] No content extracted from {original_name}")
        return {"chunks": 0, "storage_url": storage_url}

    # 3. Chunk + embed + store
    chunks = split_documents(docs)
    upsert_vectorstore(chunks, subject_id)

    text_count  = sum(1 for d in docs if d.metadata.get("type") == "text")
    image_count = sum(1 for d in docs if d.metadata.get("type") == "image_description")
    print(f"  [ingest] {original_name}: {text_count} text docs, {image_count} image descriptions → {len(chunks)} chunks")

    return {
        "chunks": len(chunks),
        "text_docs": text_count,
        "image_descriptions": image_count,
        "storage_url": storage_url
    }
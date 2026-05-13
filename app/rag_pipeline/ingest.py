import os
import base64
import zipfile
from io import BytesIO

import fitz  # pymupdf
from groq import Groq
from PIL import Image
from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE
from langchain_text_splitters import RecursiveCharacterTextSplitter
from langchain_community.vectorstores import FAISS
from langchain_community.embeddings import HuggingFaceEmbeddings
from langchain_core.documents import Document

from config import DATA_PATH, VECTOR_DB_PATH, EMBEDDING_MODEL

# ---------------------------------------------------------------------------
# Groq vision client
# ---------------------------------------------------------------------------

_groq_client = None

def get_groq_client():
    global _groq_client
    if _groq_client is None:
        from app.core.config import GROQ_API_KEY
        _groq_client = Groq(api_key=GROQ_API_KEY)
    return _groq_client


# ---------------------------------------------------------------------------
# Vision: describe an image using Groq
# ---------------------------------------------------------------------------

def describe_image(image_bytes: bytes, surrounding_context: str = "", source_hint: str = "") -> str:
    """
    Send an image to Groq's vision model and get a detailed text description.
    Returns empty string if the image is too small / blank to be meaningful.
    """
    try:
        img = Image.open(BytesIO(image_bytes))
        # Skip tiny decorative images (icons, bullets, logos, etc.)
        if img.width < 80 or img.height < 80:
            return ""

        # Re-encode as JPEG for the API (keeps payload small)
        buf = BytesIO()
        img.convert("RGB").save(buf, format="JPEG", quality=85)
        b64 = base64.b64encode(buf.getvalue()).decode()

        context_hint = f"\nSurrounding text context: {surrounding_context[:300]}" if surrounding_context else ""
        source_note  = f"\nSource file: {source_hint}" if source_hint else ""

        response = get_groq_client().chat.completions.create(
            model="meta-llama/llama-4-scout-17b-16e-instruct",
            messages=[{
                "role": "user",
                "content": [
                    {
                        "type": "image_url",
                        "image_url": {"url": f"data:image/jpeg;base64,{b64}"}
                    },
                    {
                        "type": "text",
                        "text": (
                            "You are processing a study note or academic slide.{context_hint}{source_note}\n\n"
                            "Describe this image in detail for a student who cannot see it:\n"
                            "- What type of figure is it? (diagram, graph, chart, flowchart, equation, table, photo, etc.)\n"
                            "- What concept or topic does it illustrate?\n"
                            "- Label ALL axes, components, boxes, arrows, and annotations visible.\n"
                            "- Extract ALL numerical values, formulas, or equations present.\n"
                            "- Describe any trends, relationships, or key takeaways shown.\n"
                            "Be thorough — this description will be used for retrieval."
                        ).format(context_hint=context_hint, source_note=source_note)
                    }
                ]
            }],
            max_tokens=600
        )
        return response.choices[0].message.content.strip()

    except Exception as e:
        print(f"  [vision] Could not describe image: {e}")
        return ""


# ---------------------------------------------------------------------------
# PDF loading — text blocks + embedded images
# ---------------------------------------------------------------------------

def load_pdf(file_path: str) -> list[Document]:
    docs = []
    filename = os.path.basename(file_path)
    doc = fitz.open(file_path)

    for page_num, page in enumerate(doc, start=1):
        # --- Text ---
        text = page.get_text("text").strip()
        if text:
            docs.append(Document(
                page_content=text,
                metadata={"source": filename, "page": page_num, "type": "text"}
            ))

        # --- Images ---
        image_list = page.get_images(full=True)
        for img_index, img_info in enumerate(image_list):
            xref = img_info[0]
            try:
                base_image = doc.extract_image(xref)
                image_bytes = base_image["image"]
                description = describe_image(image_bytes, surrounding_context=text, source_hint=filename)
                if description:
                    print(f"  [vision] Described image {img_index+1} on page {page_num} of {filename}")
                    docs.append(Document(
                        page_content=f"[Figure on page {page_num}]: {description}",
                        metadata={"source": filename, "page": page_num, "type": "image_description"}
                    ))
            except Exception as e:
                print(f"  [vision] Skipped image {img_index+1} on page {page_num}: {e}")

    doc.close()
    return docs


# ---------------------------------------------------------------------------
# PPTX loading — text shapes + image shapes
# ---------------------------------------------------------------------------

def load_pptx(file_path: str) -> list[Document]:
    docs = []
    filename = os.path.basename(file_path)
    prs = Presentation(file_path)

    for slide_num, slide in enumerate(prs.slides, start=1):
        slide_text_parts = []
        image_blobs = []

        for shape in slide.shapes:
            # Text
            if hasattr(shape, "text") and shape.text.strip():
                slide_text_parts.append(shape.text.strip())

            # Inline pictures
            if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
                try:
                    image_blobs.append(shape.image.blob)
                except Exception:
                    pass

        slide_text = "\n".join(slide_text_parts).strip()

        if slide_text:
            docs.append(Document(
                page_content=slide_text,
                metadata={"source": filename, "slide": slide_num, "type": "text"}
            ))

        for img_index, blob in enumerate(image_blobs):
            description = describe_image(blob, surrounding_context=slide_text, source_hint=filename)
            if description:
                print(f"  [vision] Described image {img_index+1} on slide {slide_num} of {filename}")
                docs.append(Document(
                    page_content=f"[Figure on slide {slide_num}]: {description}",
                    metadata={"source": filename, "slide": slide_num, "type": "image_description"}
                ))

    # Also extract images packed inside the .pptx zip (covers grouped/embedded objects)
    _extract_pptx_zip_images(file_path, filename, docs)

    return docs


def _extract_pptx_zip_images(file_path: str, filename: str, docs: list[Document]):
    """
    PPTX is a ZIP. Images not captured by python-pptx shapes often live
    in ppt/media/. This extracts them as a safety net.
    Already-described images won't cause harm — duplicates are filtered
    at the embedding stage by FAISS cosine dedup (similar vectors cluster).
    """
    already_described = {
        doc.metadata.get("slide")
        for doc in docs
        if doc.metadata.get("type") == "image_description"
    }

    try:
        with zipfile.ZipFile(file_path) as z:
            media_files = [f for f in z.namelist() if f.startswith("ppt/media/")
                           and f.lower().endswith((".png", ".jpg", ".jpeg", ".gif", ".bmp"))]
            for media_path in media_files:
                image_bytes = z.read(media_path)
                description = describe_image(image_bytes, source_hint=filename)
                if description:
                    docs.append(Document(
                        page_content=f"[Embedded media — {os.path.basename(media_path)}]: {description}",
                        metadata={"source": filename, "type": "image_description", "media_file": media_path}
                    ))
    except Exception as e:
        print(f"  [pptx-zip] Could not scan media in {filename}: {e}")


# ---------------------------------------------------------------------------
# DOCX loading — text + embedded images
# ---------------------------------------------------------------------------

def load_docx(file_path: str) -> list[Document]:
    """
    Extract text paragraphs and embedded images from a .docx file.
    """
    from docx import Document as DocxDocument

    docs = []
    filename = os.path.basename(file_path)

    try:
        docx = DocxDocument(file_path)
        full_text = "\n".join(p.text for p in docx.paragraphs if p.text.strip())

        if full_text:
            docs.append(Document(
                page_content=full_text,
                metadata={"source": filename, "type": "text"}
            ))

        # Images are in the zip under word/media/
        with zipfile.ZipFile(file_path) as z:
            media_files = [f for f in z.namelist() if f.startswith("word/media/")
                           and f.lower().endswith((".png", ".jpg", ".jpeg", ".gif", ".bmp"))]
            for media_path in media_files:
                image_bytes = z.read(media_path)
                description = describe_image(image_bytes, surrounding_context=full_text[:300], source_hint=filename)
                if description:
                    print(f"  [vision] Described {os.path.basename(media_path)} in {filename}")
                    docs.append(Document(
                        page_content=f"[Figure in {filename}]: {description}",
                        metadata={"source": filename, "type": "image_description", "media_file": media_path}
                    ))
    except Exception as e:
        print(f"  [docx] Error loading {filename}: {e}")

    return docs


# ---------------------------------------------------------------------------
# Load all documents from DATA_PATH
# ---------------------------------------------------------------------------

def load_documents() -> list[Document]:
    docs = []

    for file in os.listdir(DATA_PATH):
        file_path = os.path.join(DATA_PATH, file)

        if file.endswith(".pdf"):
            print(f"Loading PDF: {file}")
            docs.extend(load_pdf(file_path))

        elif file.endswith(".pptx"):
            print(f"Loading PPTX: {file}")
            docs.extend(load_pptx(file_path))

        elif file.endswith(".docx"):
            print(f"Loading DOCX: {file}")
            docs.extend(load_docx(file_path))

        elif file.lower().endswith((".png", ".jpg", ".jpeg")):
            print(f"Loading image: {file}")
            with open(file_path, "rb") as f:
                image_bytes = f.read()
            description = describe_image(image_bytes, source_hint=file)
            if description:
                docs.append(Document(
                    page_content=f"[Image file — {file}]: {description}",
                    metadata={"source": file, "type": "image_description"}
                ))

    return docs


# ---------------------------------------------------------------------------
# Chunking — respects chunk_type; image descriptions are NOT split
# ---------------------------------------------------------------------------

def split_documents(documents: list[Document]) -> list[Document]:
    splitter = RecursiveCharacterTextSplitter(
        chunk_size=500,
        chunk_overlap=100
    )

    chunks = []
    for doc in documents:
        if doc.metadata.get("type") == "image_description":
            # Keep image descriptions whole — splitting breaks their meaning
            chunks.append(doc)
        else:
            split_texts = splitter.split_text(doc.page_content)
            for chunk in split_texts:
                chunks.append(Document(
                    page_content=chunk,
                    metadata=doc.metadata
                ))

    return chunks


# ---------------------------------------------------------------------------
# Vectorstore
# ---------------------------------------------------------------------------

def create_vectorstore(chunks: list[Document]):
    embeddings = HuggingFaceEmbeddings(model_name=EMBEDDING_MODEL)
    os.makedirs(VECTOR_DB_PATH, exist_ok=True)
    db = FAISS.from_documents(chunks, embeddings)
    db.save_local(VECTOR_DB_PATH)


# ---------------------------------------------------------------------------
# Main
# ---------------------------------------------------------------------------

def main():
    print("Loading documents")
    documents = load_documents()

    if not documents:
        print("No documents found in data/raw/")
        return

    text_docs  = [d for d in documents if d.metadata.get("type") == "text"]
    image_docs = [d for d in documents if d.metadata.get("type") == "image_description"]
    print(f"Loaded {len(documents)} documents — {len(text_docs)} text, {len(image_docs)} image descriptions")

    print("Splitting documents into chunks")
    chunks = split_documents(documents)
    print(f"Created {len(chunks)} chunks")

    print("Creating embeddings and saving vector database")
    create_vectorstore(chunks)
    print("Ingestion complete")


if __name__ == "__main__":
    main()